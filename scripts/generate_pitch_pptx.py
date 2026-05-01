import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(md_file_path, output_pptx_path):
    prs = Presentation()

    # Define some colors
    # Corridor brand colors (assumed based on "Financial OS" theme - dark, premium)
    PRIMARY_COLOR = RGBColor(0, 51, 102) # Dark Blue
    ACCENT_COLOR = RGBColor(0, 153, 204) # Lighter Blue
    TEXT_COLOR = RGBColor(50, 50, 50) # Dark Gray

    with open(md_file_path, 'r') as f:
        lines = f.readlines()

    current_slide = None
    current_body = None
    
    # Slide Layouts
    # 0: Title Slide
    # 1: Title and Content
    # 2: Section Header
    
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if line.startswith('# '):
            # Title Slide
            slide_layout = prs.slide_layouts[0]
            current_slide = prs.slides.add_slide(slide_layout)
            title = current_slide.shapes.title
            subtitle = current_slide.placeholders[1]
            
            title.text = line[2:].strip()
            subtitle.text = "Financial Operating System" # Default subtitle or extract from next lines

        elif line.startswith('## '):
            # New Content Slide
            slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(slide_layout)
            title = current_slide.shapes.title
            title.text = line[3:].strip()
            
            # Reset body for new slide
            current_body = current_slide.placeholders[1].text_frame
            current_body.clear() 

        elif line.startswith('### '):
            # Sub-section within a slide (bold text or new paragraph)
            if current_body:
                p = current_body.add_paragraph()
                p.text = line[4:].strip()
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = PRIMARY_COLOR
                p.space_before = Pt(12)

        elif line.startswith('* ') or line.startswith('- '):
            # Bullet point
            if current_body:
                p = current_body.add_paragraph()
                text = line[2:].strip()
                
                # Handle bolding in markdown **text**
                parts = text.split('**')
                for i, part in enumerate(parts):
                    run = p.add_run()
                    run.text = part
                    if i % 2 == 1: # Odd parts are between **
                        run.font.bold = True
                
                p.level = 0
                p.space_after = Pt(6)

        else:
            # Normal text
            if current_body:
                # If it's the first paragraph, use the existing one, else add new
                if len(current_body.paragraphs) == 0:
                    p = current_body.add_paragraph()
                else:
                    p = current_body.add_paragraph()
                
                p.text = line
                p.level = 0
                p.font.size = Pt(18)

    # Save
    prs.save(output_pptx_path)
    print(f"Presentation saved to {output_pptx_path}")

if __name__ == "__main__":
    md_path = os.path.join(os.path.dirname(__file__), '../docs/pitch_document.md')
    pptx_path = os.path.join(os.path.dirname(__file__), '../docs/Corridor_Pitch.pptx')
    
    # Ensure absolute paths
    md_path = os.path.abspath(md_path)
    pptx_path = os.path.abspath(pptx_path)
    
    print(f"Reading from: {md_path}")
    create_presentation(md_path, pptx_path)
